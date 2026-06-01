"""
Inserts a "Where We Are" slide at position 0 in lm_lead_time_timeline.pptx
Sources:
  • W21 Confluence weekly report + delivery table (CW14-CW20 actuals)
  • WS6 main Confluence page (1505067229) – goals and milestones
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree
import copy

PPTX_PATH = "/Users/ana.lira/Documents/Orbis/Repos/Scaling LM Releases Expert/lm_lead_time_timeline.pptx"

prs = Presentation(PPTX_PATH)
W = prs.slide_width.inches   # 13.33
H = prs.slide_height.inches  # 7.5

# ── Colours (matching existing slide) ────────────────────────────────────────
C_BG     = RGBColor(0x1A, 0x1A, 0x2E)
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_DIM    = RGBColor(0x8A, 0x8A, 0xAA)
C_PURPLE = RGBColor(0x7C, 0x4D, 0xFF)
C_GREEN  = RGBColor(0x00, 0xC9, 0x8D)
C_ORANGE = RGBColor(0xFF, 0x8C, 0x00)
C_BLUE   = RGBColor(0x4F, 0xB3, 0xFF)
C_GOLD   = RGBColor(0xFF, 0xDD, 0x55)
C_DARK   = RGBColor(0x11, 0x11, 0x22)
C_CARD   = RGBColor(0x24, 0x24, 0x42)   # slightly lighter than bg for cards
C_STAR   = RGBColor(0xFF, 0xEE, 0x00)

# ── Add blank slide ───────────────────────────────────────────────────────────
blank_layout = prs.slide_layouts[6]
new_slide = prs.slides.add_slide(blank_layout)

# Move new slide to position 0 (insert before existing slide)
xml_slides = prs.slides._sldIdLst
# The new slide is currently last; move it to front
last = xml_slides[-1]
xml_slides.remove(last)
xml_slides.insert(0, last)

slide = new_slide

# ── Background ────────────────────────────────────────────────────────────────
bg = slide.background
fill = bg.fill
fill.solid()
fill.fore_color.rgb = C_BG

# ── Helpers ───────────────────────────────────────────────────────────────────
def txt(slide, text, x, y, w, h, size=11, bold=False, color=None,
        align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color or C_WHITE
    return tb

def rect(slide, x, y, w, h, fill_color, line_color=None, line_pt=None, radius=False):
    from pptx.util import Pt as ptU
    sh = slide.shapes.add_shape(
        1,
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill_color
    if line_color:
        sh.line.color.rgb = line_color
        if line_pt:
            sh.line.width = Pt(line_pt)
    else:
        sh.line.fill.background()
    return sh

def hline(slide, x1, x2, y, color, pt=1.5):
    c = slide.shapes.add_connector(1, Inches(x1), Inches(y), Inches(x2), Inches(y))
    c.line.color.rgb = color
    c.line.width = Pt(pt)
    return c

def vline(slide, x, y1, y2, color, pt=1.0):
    c = slide.shapes.add_connector(1, Inches(x), Inches(y1), Inches(x), Inches(y2))
    c.line.color.rgb = color
    c.line.width = Pt(pt)
    return c

# ═══════════════════════════════════════════════════════════════════════════════
# TITLE
# ═══════════════════════════════════════════════════════════════════════════════
txt(slide, "WS6 – Where We Are", 0.30, 0.12, 8.0, 0.60,
    size=28, bold=True, color=C_WHITE)
txt(slide, "Enabling Weekly LM Map Releases  ·  March → May 2026",
    0.30, 0.68, 9.0, 0.35, size=12, color=C_DIM)

hline(slide, 0.25, W - 0.25, 1.05, C_DIM, pt=0.6)

# ═══════════════════════════════════════════════════════════════════════════════
# LEFT COLUMN — Achievements
# ═══════════════════════════════════════════════════════════════════════════════
COL1_X = 0.30
COL1_W = 4.60

txt(slide, "What we've delivered", COL1_X, 1.15, COL1_W, 0.38,
    size=13, bold=True, color=C_BLUE)

achievements = [
    ("✅", "9 weekly releases", "CW14–CW21, every week without a gap"),
    ("✅", "7 of 9 on time",   "On or before Thursday target"),
    ("✅", "First combined CARIAD\n    + VOLVO release", "CW17 — both endpoints live same week"),
    ("✅", "CARIAD confidence restored", "Voluntarily skipped CW19 — a sign of trust"),
    ("✅", "New features shipped", "Drive Path layer + 8 self-contained features (CW18)"),
    ("✅", "Incremental pipeline live", "New mini-binding layer active; old one decommissioned"),
]

ay = 1.58
for icon, headline, sub in achievements:
    txt(slide, icon, COL1_X, ay, 0.35, 0.28, size=12)
    txt(slide, headline, COL1_X + 0.35, ay - 0.01, COL1_W - 0.40, 0.30,
        size=11, bold=True, color=C_WHITE)
    txt(slide, sub, COL1_X + 0.35, ay + 0.25, COL1_W - 0.40, 0.28,
        size=9, color=C_DIM, italic=True)
    ay += 0.62

# Separator
vline(slide, COL1_X + COL1_W + 0.20, 1.10, 7.30, C_CARD, pt=1.2)

# ═══════════════════════════════════════════════════════════════════════════════
# CENTER COLUMN — Lead Time Goals Ladder
# ═══════════════════════════════════════════════════════════════════════════════
COL2_X = 5.40
COL2_W = 3.80

txt(slide, "Lead time targets", COL2_X, 1.15, COL2_W, 0.38,
    size=13, bold=True, color=C_BLUE)

# Current marker
rect(slide, COL2_X, 1.58, COL2_W, 0.62, RGBColor(0x2A, 0x2A, 0x4A))
txt(slide, "NOW", COL2_X + 0.12, 1.61, 0.60, 0.22, size=8, bold=True, color=C_DIM)
txt(slide, "~13.8 days", COL2_X + 0.12, 1.82, 1.60, 0.28,
    size=15, bold=True, color=C_WHITE)
txt(slide, "avg CW16–CW20", COL2_X + 1.80, 1.90, 1.80, 0.22, size=8.5, color=C_DIM)

# Arrow down
vline(slide, COL2_X + COL2_W/2, 2.22, 2.40, C_DIM, pt=1.0)
txt(slide, "▼", COL2_X + COL2_W/2 - 0.08, 2.34, 0.25, 0.22, size=9, color=C_DIM)

# Goal 1 — Q2 end
rect(slide, COL2_X, 2.52, COL2_W, 0.72, RGBColor(0x28, 0x1A, 0x55))
hline(slide, COL2_X, COL2_X + COL2_W, 2.52, C_PURPLE, pt=3)
txt(slide, "End Q2  ·  External goal", COL2_X + 0.12, 2.56, COL2_W - 0.20, 0.22,
    size=8.5, bold=True, color=C_PURPLE)
txt(slide, "≤ 11.5 days", COL2_X + 0.12, 2.78, 2.20, 0.36,
    size=18, bold=True, color=C_PURPLE)
# Star badge — already done once
rect(slide, COL2_X + 2.50, 2.60, 1.15, 0.54, RGBColor(0x3A, 0x2A, 0x10))
txt(slide, "⭐ Already hit it!", COL2_X + 2.54, 2.63, 1.10, 0.22,
    size=9, bold=True, color=C_GOLD)
txt(slide, "CW17: 10.61 days", COL2_X + 2.54, 2.84, 1.10, 0.22,
    size=8.5, color=C_GOLD)

# Arrow down
vline(slide, COL2_X + COL2_W/2, 3.26, 3.44, C_DIM, pt=1.0)
txt(slide, "▼", COL2_X + COL2_W/2 - 0.08, 3.38, 0.25, 0.22, size=9, color=C_DIM)

# Goal 2 — September Q3
rect(slide, COL2_X, 3.56, COL2_W, 0.72, RGBColor(0x10, 0x2E, 0x26))
hline(slide, COL2_X, COL2_X + COL2_W, 3.56, C_GREEN, pt=3)
txt(slide, "September  ·  External goal", COL2_X + 0.12, 3.60, COL2_W - 0.20, 0.22,
    size=8.5, bold=True, color=C_GREEN)
txt(slide, "≤ 9 days", COL2_X + 0.12, 3.82, 2.20, 0.36,
    size=18, bold=True, color=C_GREEN)
txt(slide, "Incremental\npipeline", COL2_X + 2.80, 3.72, 0.85, 0.42,
    size=8, color=C_GREEN, italic=True)

# Arrow down
vline(slide, COL2_X + COL2_W/2, 4.30, 4.48, C_DIM, pt=1.0)
txt(slide, "▼", COL2_X + COL2_W/2 - 0.08, 4.42, 0.25, 0.22, size=9, color=C_DIM)

# Goal 3 — EOY internal
rect(slide, COL2_X, 4.60, COL2_W, 0.72, RGBColor(0x2B, 0x1A, 0x00))
hline(slide, COL2_X, COL2_X + COL2_W, 4.60, C_ORANGE, pt=3)
txt(slide, "End of Year  ·  Internal goal", COL2_X + 0.12, 4.64, COL2_W - 0.20, 0.22,
    size=8.5, bold=True, color=C_ORANGE)
txt(slide, "≤ 7 days", COL2_X + 0.12, 4.86, 2.20, 0.36,
    size=18, bold=True, color=C_ORANGE)
txt(slide, "Ready for\nCARIAD mid-'27", COL2_X + 2.70, 4.76, 0.95, 0.42,
    size=8, color=C_ORANGE, italic=True)

# Arrow down
vline(slide, COL2_X + COL2_W/2, 5.34, 5.52, C_DIM, pt=1.0)
txt(slide, "▼", COL2_X + COL2_W/2 - 0.08, 5.46, 0.25, 0.22, size=9, color=C_DIM)

# North star
rect(slide, COL2_X, 5.64, COL2_W, 0.58, RGBColor(0x0D, 0x0D, 0x20))
hline(slide, COL2_X, COL2_X + COL2_W, 5.64, C_BLUE, pt=1)
txt(slide, "🎯  North Star  ·  mid-2027", COL2_X + 0.12, 5.68, COL2_W - 0.20, 0.22,
    size=8.5, bold=True, color=C_BLUE)
txt(slide, "CARIAD weekly cadence at scale", COL2_X + 0.12, 5.90, COL2_W - 0.20, 0.26,
    size=10, color=C_BLUE)

# ═══════════════════════════════════════════════════════════════════════════════
# RIGHT COLUMN — Momentum callout
# ═══════════════════════════════════════════════════════════════════════════════
vline(slide, COL2_X + COL2_W + 0.20, 1.10, 7.30, C_CARD, pt=1.2)

COL3_X = COL2_X + COL2_W + 0.50
COL3_W = W - COL3_X - 0.20

txt(slide, "Momentum", COL3_X, 1.15, COL3_W, 0.38,
    size=13, bold=True, color=C_BLUE)

# Mini lead-time bar chart
bars = [
    ("CW15", 13.82, C_DIM),
    ("CW16", 14.44, C_DIM),
    ("CW17", 10.61, C_GOLD),   # best
    ("CW18", 13.81, C_DIM),
    ("CW20", 13.77, C_BLUE),   # latest
]
BAR_X    = COL3_X + 0.10
BAR_AREA_W = COL3_W - 0.20
MAX_LT   = 16.0
BAR_H    = 0.28
BAR_GAP  = 0.10
chart_y  = 1.60

txt(slide, "Lead time per release (days)", BAR_X, chart_y - 0.02, BAR_AREA_W, 0.22,
    size=8, color=C_DIM, italic=True)
chart_y += 0.24

for label, lt, color in bars:
    bar_w = BAR_AREA_W * lt / MAX_LT
    rect(slide, BAR_X, chart_y, bar_w, BAR_H, color)
    txt(slide, f"{lt}", BAR_X + bar_w + 0.04, chart_y + 0.02, 0.55, BAR_H,
        size=9, bold=(lt < 11), color=color)
    txt(slide, label, BAR_X - 0.38, chart_y + 0.02, 0.36, BAR_H,
        size=8.5, color=C_DIM, align=PP_ALIGN.RIGHT)
    chart_y += BAR_H + BAR_GAP

# Target line (11.5d)
target_x = BAR_X + BAR_AREA_W * 11.5 / MAX_LT
vline(slide, target_x, 1.84, chart_y - BAR_GAP + 0.06, C_PURPLE, pt=1.2)
txt(slide, "Q2 target", target_x - 0.60, chart_y - BAR_GAP + 0.06, 0.75, 0.22,
    size=7.5, color=C_PURPLE, align=PP_ALIGN.RIGHT)

# CW17 callout box
chart_y += 0.12
rect(slide, COL3_X, chart_y, COL3_W - 0.05, 1.00, RGBColor(0x1E, 0x1E, 0x10))
hline(slide, COL3_X, COL3_X + COL3_W - 0.05, chart_y, C_GOLD, pt=2.5)
txt(slide, "⭐  CW17: 10.61 days", COL3_X + 0.12, chart_y + 0.06, COL3_W - 0.25, 0.30,
    size=13, bold=True, color=C_GOLD)
txt(slide, "First time under Q2 target.",
    COL3_X + 0.12, chart_y + 0.38, COL3_W - 0.25, 0.22,
    size=9.5, color=C_WHITE)
txt(slide, "Proof that 11.5 days is achievable now.",
    COL3_X + 0.12, chart_y + 0.60, COL3_W - 0.25, 0.28,
    size=9.5, bold=True, color=C_GOLD)

# ── Footer ────────────────────────────────────────────────────────────────────
hline(slide, 0.25, W - 0.25, 7.10, C_DIM, pt=0.4)
txt(slide, "Sources: W21 Confluence weekly report · WS6 main page (1505067229) · JIRA ADASHD-2672  ·  May 2026",
    0.25, 7.15, W - 0.50, 0.25, size=7.5, color=C_DIM)

# ── Save ──────────────────────────────────────────────────────────────────────
prs.save(PPTX_PATH)
print(f"Saved: {PPTX_PATH}")
print(f"Slides: {len(prs.slides)}")
