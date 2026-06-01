"""
WS6 Lead Time Improvement Timeline – PowerPoint slide
Sources:
  • JIRA ADASHD-3339, ADASHD-3625, ADASHD-3336, ADASHD-2756, ADASHD-2755, ADASHD-2764
  • Confluence WS6 main page (1505067229) – CW20 actuals: 13.77 days
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import datetime

# ── Slide dimensions (widescreen 16:9) ──────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

slide_layout = prs.slide_layouts[6]  # blank
slide = prs.slides.add_slide(slide_layout)

# ── Colours ──────────────────────────────────────────────────────────────────
C_BG        = RGBColor(0x1A, 0x1A, 0x2E)   # dark navy bg
C_TITLE     = RGBColor(0xFF, 0xFF, 0xFF)
C_AXIS      = RGBColor(0x8A, 0x8A, 0xAA)
C_PURPLE    = RGBColor(0x7C, 0x4D, 0xFF)   # process / automation
C_GREEN     = RGBColor(0x00, 0xC9, 0x8D)   # incremental
C_ORANGE    = RGBColor(0xFF, 0x8C, 0x00)   # target lines
C_RED_LINE  = RGBColor(0xFF, 0x45, 0x45)   # current
C_STEP      = RGBColor(0x4F, 0xB3, 0xFF)   # staircase line
C_TICK      = RGBColor(0xAA, 0xAA, 0xCC)
C_LABEL_DRK = RGBColor(0x11, 0x11, 0x22)

# ── Background ───────────────────────────────────────────────────────────────
from pptx.oxml.ns import qn
from lxml import etree

bg = slide.background
fill = bg.fill
fill.solid()
fill.fore_color.rgb = C_BG

# ── Helper: add text box ──────────────────────────────────────────────────────
def add_text(slide, text, x, y, w, h, size=10, bold=False, color=None,
             align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color or C_TITLE
    return txb

# ── Helper: filled rectangle ─────────────────────────────────────────────────
def add_rect(slide, x, y, w, h, fill_color, line_color=None, line_w=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        if line_w:
            shape.line.width = line_w
    else:
        shape.line.fill.background()
    return shape

# ── Helper: horizontal line ───────────────────────────────────────────────────
def add_hline(slide, x1, x2, y, color, width_pt=1.5):
    from pptx.util import Pt as ptU
    connector = slide.shapes.add_connector(
        1,  # straight
        Inches(x1), Inches(y), Inches(x2), Inches(y)
    )
    connector.line.color.rgb = color
    connector.line.width = Pt(width_pt)
    return connector

def add_vline(slide, x, y1, y2, color, width_pt=1.0):
    connector = slide.shapes.add_connector(
        1,
        Inches(x), Inches(y1), Inches(x), Inches(y2)
    )
    connector.line.color.rgb = color
    connector.line.width = Pt(width_pt)
    return connector

# ═══════════════════════════════════════════════════════════════════════════════
# LAYOUT CONSTANTS
# ═══════════════════════════════════════════════════════════════════════════════
MARGIN_L = 1.1   # left margin where timeline starts
MARGIN_R = 0.25  # right margin
TL_X1    = MARGIN_L
TL_X2    = prs.slide_width.inches - MARGIN_R
TL_W     = TL_X2 - TL_X1

# Timeline spans May 1 → Sep 30  (153 days)
import datetime as dt
T_START = dt.date(2026, 5, 1)
T_END   = dt.date(2026, 9, 30)
T_DAYS  = (T_END - T_START).days

def date_x(d):
    """Return x-inch position for a date."""
    delta = (d - T_START).days
    return TL_X1 + TL_W * delta / T_DAYS

# Month tick positions
MONTHS = [
    (dt.date(2026, 5,  1), "May '26"),
    (dt.date(2026, 6,  1), "Jun"),
    (dt.date(2026, 7,  1), "Jul"),
    (dt.date(2026, 8,  1), "Aug"),
    (dt.date(2026, 9,  1), "Sep"),
    (dt.date(2026, 9, 30), ""),
]

TODAY       = dt.date(2026, 5, 21)
TODAY_X     = date_x(TODAY)

# Gantt section
GANTT_Y_TOP = 1.15
GANTT_ROW_H = 0.36
GANTT_BAR_H = 0.22
GANTT_N     = 8

# Chart section
CHART_Y_TOP  = GANTT_Y_TOP + GANTT_N * GANTT_ROW_H + 0.30
CHART_Y_BOT  = 7.15
CHART_H      = CHART_Y_BOT - CHART_Y_TOP

LT_MAX = 15.0
LT_MIN = 6.5

def lt_y(lt_val):
    """Y position for a lead-time value."""
    frac = (lt_val - LT_MIN) / (LT_MAX - LT_MIN)
    return CHART_Y_BOT - frac * CHART_H

# ═══════════════════════════════════════════════════════════════════════════════
# TITLE
# ═══════════════════════════════════════════════════════════════════════════════
add_text(slide, "WS6 – Lead Time Improvement Roadmap  (Goal 2)",
         0.15, 0.10, 10.0, 0.55, size=18, bold=True, color=C_TITLE)
add_text(slide, "May → September 2026  |  Current: 13.77 days (CW20)  |  Target: 9.0 days (Q3)",
         0.15, 0.62, 11.0, 0.35, size=11, color=C_AXIS)

# ═══════════════════════════════════════════════════════════════════════════════
# MONTH AXIS (shared)
# ═══════════════════════════════════════════════════════════════════════════════
AXIS_Y = GANTT_Y_TOP - 0.22
add_hline(slide, TL_X1, TL_X2, AXIS_Y + 0.12, C_AXIS, width_pt=0.8)

for d, label in MONTHS:
    x = date_x(d)
    add_vline(slide, x, AXIS_Y + 0.08, AXIS_Y + 0.16, C_AXIS, width_pt=0.8)
    if label:
        add_text(slide, label, x - 0.18, AXIS_Y - 0.12, 0.6, 0.22,
                 size=9, color=C_AXIS, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════════
# GANTT ROWS
# ═══════════════════════════════════════════════════════════════════════════════
initiatives = [
    # (label, ticket, owner, start_date, end_date, color, saving_text)
    ("To-Be process design",         "ADASHD-2756", "M. Aigner",
     dt.date(2026, 5, 1), dt.date(2026, 5, 22),   C_PURPLE, "Alignment"),
    ("Process telemetry",            "ADASHD-3336", "A. Jalali",
     dt.date(2026, 5, 1), dt.date(2026, 5, 29),   C_PURPLE, "Enables measurement"),
    ("Remove NDS.Live idle time",    "ADASHD-3339", "P. van Hulten",
     dt.date(2026, 5, 1), dt.date(2026, 5, 31),   C_PURPLE, "−2.3 days"),
    ("Tool automation / handoffs",   "ADASHD-3625", "A. Jalali",
     dt.date(2026, 5, 22), dt.date(2026, 6, 21),  C_PURPLE, "−0.5 days"),
    ("DevMap/Staging stability",     "ADASHD-2755", "T. Szambor",
     dt.date(2026, 5, 15), dt.date(2026, 6, 30),  C_PURPLE, "−0.5 days"),
    ("Incremental P0a: Crosslink",   "ADASHD-2764", "T. Bekaert",
     dt.date(2026, 6, 15), dt.date(2026, 8, 1),   C_GREEN,  "−1.2 days"),
    ("Incremental P0b: Lane Fn",     "ADASHD-2764", "E. Goraca",
     dt.date(2026, 7, 1),  dt.date(2026, 9, 1),   C_GREEN,  "−2.0 days"),
    ("Incremental P1: Self-contained","ADASHD-2764", "D. Turobos",
     dt.date(2026, 7, 15), dt.date(2026, 9, 15),  C_GREEN,  "−0.8 days"),
]

LABEL_W = MARGIN_L - 0.10

for i, (label, ticket, owner, d_start, d_end, color, saving) in enumerate(initiatives):
    row_y  = GANTT_Y_TOP + i * GANTT_ROW_H
    bar_y  = row_y + (GANTT_ROW_H - GANTT_BAR_H) / 2

    # row separator (subtle)
    sep_col = RGBColor(0x2A, 0x2A, 0x4A)
    add_hline(slide, 0.05, TL_X2, row_y, sep_col, width_pt=0.4)

    # left label
    add_text(slide, label, 0.08, bar_y - 0.03, LABEL_W - 0.30, 0.26,
             size=9, bold=False, color=C_TITLE)
    add_text(slide, ticket, 0.08, bar_y + 0.18, LABEL_W - 0.30, 0.18,
             size=7.5, color=C_AXIS)

    # bar
    bx = date_x(d_start)
    bw = date_x(d_end) - bx
    if bw < 0.05:
        bw = 0.05
    bar = add_rect(slide, bx, bar_y, bw, GANTT_BAR_H, color)

    # saving label inside or to the right of bar
    save_x = bx + bw + 0.04
    if save_x + 0.7 > TL_X2:
        save_x = bx + 0.04
        save_col = C_LABEL_DRK
    else:
        save_col = color
    add_text(slide, saving, save_x, bar_y + 0.01, 1.0, 0.22,
             size=8.5, bold=True, color=save_col)

    # owner
    add_text(slide, owner, save_x, bar_y + 0.20, 1.1, 0.18,
             size=7.5, color=C_AXIS)

# ═══════════════════════════════════════════════════════════════════════════════
# TODAY vertical line (shared across both sections)
# ═══════════════════════════════════════════════════════════════════════════════
add_vline(slide, TODAY_X, AXIS_Y + 0.10,
          CHART_Y_BOT, RGBColor(0xFF, 0xDD, 0x55), width_pt=1.2)
add_text(slide, "Today", TODAY_X - 0.22, GANTT_Y_TOP - 0.05, 0.55, 0.18,
         size=8, bold=True, color=RGBColor(0xFF, 0xDD, 0x55), align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════════
# LEAD TIME STAIRCASE CHART
# ═══════════════════════════════════════════════════════════════════════════════
# Y-axis labels
for lt_val in [7, 8, 9, 10, 11, 12, 13, 14, 15]:
    y = lt_y(lt_val)
    add_text(slide, str(lt_val), TL_X1 - 0.40, y - 0.12, 0.30, 0.25,
             size=8, color=C_AXIS, align=PP_ALIGN.RIGHT)
    add_hline(slide, TL_X1 - 0.04, TL_X2, y,
              RGBColor(0x2A, 0x2A, 0x45), width_pt=0.5)

# Y-axis label
add_text(slide, "Lead time (days)", 0.02, CHART_Y_TOP + CHART_H / 2 - 0.15,
         0.9, 0.30, size=9, color=C_AXIS)

# Staircase steps: (date, lead_time_after, annotation)
steps = [
    (dt.date(2026, 5,  21), 13.77, "Now 13.77d"),
    (dt.date(2026, 5,  31), 11.47, "−2.3d\n(NDS.Live idle)"),
    (dt.date(2026, 6,  30), 10.47, "−1.0d\n(automation\n+ stability)"),
    (dt.date(2026, 8,   1),  9.27, "−1.2d\n(Crosslink\nincremental)"),
    (dt.date(2026, 9,   1),  7.27, "−2.0d\n(Lane Fn\nincremental)"),
    (dt.date(2026, 9,  30),  6.47, "−0.8d\n(P1 incremental)"),
]

# Target lines
Q2_LT = 11.5
Q3_LT = 9.0

add_hline(slide, TL_X1, TL_X2, lt_y(Q2_LT), C_ORANGE, width_pt=1.2)
add_text(slide, f"Q2 target  {Q2_LT}d",
         TL_X2 + 0.05, lt_y(Q2_LT) - 0.12, 1.1, 0.25,
         size=8.5, bold=True, color=C_ORANGE)

add_hline(slide, TL_X1, TL_X2, lt_y(Q3_LT), C_GREEN, width_pt=1.2)
add_text(slide, f"Q3 target  {Q3_LT}d",
         TL_X2 + 0.05, lt_y(Q3_LT) - 0.12, 1.1, 0.25,
         size=8.5, bold=True, color=C_GREEN)

# Draw staircase as a series of horizontal + vertical lines
for i in range(len(steps) - 1):
    d0, lt0, _ = steps[i]
    d1, lt1, _ = steps[i + 1]
    x0, y0 = date_x(d0), lt_y(lt0)
    x1, y1 = date_x(d1), lt_y(lt1)
    # horizontal segment from d0 to d1 at lt0
    add_hline(slide, x0, x1, y0, C_STEP, width_pt=2.5)
    # vertical drop at d1
    add_vline(slide, x1, y0, y1, C_STEP, width_pt=2.5)

# Extend last step to end
last_x = date_x(steps[-1][0])
last_y = lt_y(steps[-1][1])
add_hline(slide, last_x, TL_X2, last_y, C_STEP, width_pt=2.5)

# Dot + annotation at each step change
for i, (d, lt, note) in enumerate(steps):
    x = date_x(d)
    y = lt_y(lt)
    dot = add_rect(slide, x - 0.06, y - 0.06, 0.12, 0.12, C_STEP)
    if note:
        ann_y = y - 0.65 if lt < 13 else y + 0.08
        ann_x = x - 0.55
        if i == 0:
            ann_x = x + 0.08
        add_text(slide, note, ann_x, ann_y, 1.1, 0.65,
                 size=8, color=C_STEP, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════════
# LEGEND
# ═══════════════════════════════════════════════════════════════════════════════
LEG_X = 0.08
LEG_Y = 7.00
add_rect(slide, LEG_X, LEG_Y, 0.25, 0.15, C_PURPLE)
add_text(slide, "Process / automation", LEG_X + 0.30, LEG_Y - 0.02, 1.6, 0.20,
         size=8.5, color=C_TITLE)
add_rect(slide, LEG_X + 2.1, LEG_Y, 0.25, 0.15, C_GREEN)
add_text(slide, "Incremental processing", LEG_X + 2.40, LEG_Y - 0.02, 1.7, 0.20,
         size=8.5, color=C_TITLE)

# ═══════════════════════════════════════════════════════════════════════════════
# SAVE
# ═══════════════════════════════════════════════════════════════════════════════
OUT = "/Users/ana.lira/Documents/Orbis/Repos/Scaling LM Releases Expert/lm_lead_time_timeline.pptx"
prs.save(OUT)
print(f"Saved: {OUT}")
